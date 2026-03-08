"""Document Ingestion Pipeline — chunk, embed, and store in pgvector.

Supported formats:
    .txt   — Plain text files
    .pdf   — PDF documents (with OCR fallback for scanned pages)
    .docx  — Microsoft Word documents
    .xlsx  — Excel spreadsheets (openpyxl)
    .xls   — Legacy Excel spreadsheets (xlrd)
    .pptx  — PowerPoint presentations (python-pptx)
    .md    — Markdown files
    .rtf   — Rich Text Format (striprtf)
    .doc   — Legacy Word binary (best-effort via olefile)
    .xml   — XML documents
    .csv   — CSV files (generic + Purview eDiscovery)
    .tsv   — Tab-separated values
    .log   — Log files (plain text)
    .zip   — Archive files (recursive extraction with safety guards)
    .json  — JSON files (generic + Purview eDiscovery)
    .html / .htm — HTML files
    .png / .jpg / .jpeg / .tiff / .bmp — Images (via OCR)
    .eml   — Standard email files (with attachment processing)
    .msg   — Outlook message files (with attachment processing)
    .pst   — Outlook data files (requires Outlook on Windows)

Usage:
    python -m app.ingest --dir ./data/raw_corpus
    python -m app.ingest --dir ./data/mixed_docs --doc-type "corporate-comms"
"""

from __future__ import annotations

import argparse
import base64
import email
import io
import json
import logging
import os
import platform
import re
import tempfile
from email import policy
from html import unescape
from pathlib import Path

from langchain_core.documents import Document
from pypdf import PdfReader

from app.blob_storage import is_blob_mode, iter_files_for_source
from app.config import settings
from app.db import init_db, upsert_chunks, upsert_pst_email_metadata

logger = logging.getLogger(__name__)


# ── Ingestion report for auditing success/failure ─────────────────────

from dataclasses import dataclass, field as _field


@dataclass
class IngestionReport:
    """Tracks success/failure counts during document ingestion."""

    total_files_found: int = 0
    files_loaded: int = 0
    files_failed: int = 0
    failed_files: list[dict[str, str]] = _field(default_factory=list)
    chunks_created: int = 0
    chunks_stored: int = 0

    def record_failure(self, filename: str, error: str) -> None:
        self.files_failed += 1
        self.failed_files.append({"file": filename, "error": error[:200]})

    def summary(self) -> str:
        status = "OK" if self.files_failed == 0 else "PARTIAL"
        return (
            f"[{status}] files_found={self.total_files_found} "
            f"loaded={self.files_loaded} failed={self.files_failed} "
            f"chunks={self.chunks_stored}"
        )


# File extensions handled by each loader
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"}
_IMAGE_MIMES = {"image/png", "image/jpeg", "image/jpg", "image/gif", "image/bmp",
                "image/tiff", "image/webp"}
_AUDIO_MIMES = {"audio/mpeg", "audio/mp3", "audio/wav", "audio/x-wav", "audio/ogg",
                "audio/flac", "audio/mp4", "audio/m4a", "audio/webm", "audio/x-m4a"}
_AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".webm"}
_ALL_SUPPORTED = {
    ".txt",
    ".pdf",
    ".docx",
    ".xlsx", ".xls",
    ".pptx",
    ".md",
    ".rtf", ".doc",
    ".xml",
    ".json",
    ".csv", ".tsv",
    ".log",
    ".zip",
    ".html", ".htm",
    *_IMAGE_EXTS,
    ".eml",
    ".msg",
    ".pst",
}

# ── ZIP archive safety limits ──────────────────────────────────────
_ZIP_MAX_NESTING_DEPTH = 2
_ZIP_MAX_TOTAL_BYTES = 500 * 1024 * 1024   # 500 MB
_ZIP_MAX_FILE_COUNT = 10_000

# ── JSON flattening depth limit ────────────────────────────────────
_JSON_MAX_DEPTH = 50

# ── CSV/TSV batching ──────────────────────────────────────────────
_CSV_BATCH_SIZE = 50


# ─────────────────────────────────────────────────────────────────────
#  OCR helper (shared by image loader + scanned-PDF fallback)
# ─────────────────────────────────────────────────────────────────────
def _ocr_image(image) -> str:
    """Run Tesseract OCR on a PIL Image. Returns extracted text or ''."""
    try:
        import pytesseract
        text = pytesseract.image_to_string(image, lang="eng")
        return text.strip()
    except Exception as exc:
        logger.warning("OCR failed: %s (is Tesseract installed?)", exc)
        return ""


# ─────────────────────────────────────────────────────────────────────
#  Attachment processing — GPT-4o Vision OCR + Whisper transcription
# ─────────────────────────────────────────────────────────────────────
def _ocr_image_gpt4o(image_bytes: bytes, filename: str = "image") -> str:
    """Extract text from an image using Azure GPT-4o vision.

    Sends the image as a base64 data-URI to GPT-4o and asks it to
    extract all visible text/content.  Falls back to Tesseract if
    the API call fails.
    """
    try:
        from openai import AzureOpenAI

        # Detect MIME type from extension or default to png
        ext = Path(filename).suffix.lower()
        mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                    ".gif": "image/gif", ".bmp": "image/bmp", ".tiff": "image/tiff",
                    ".webp": "image/webp"}
        mime = mime_map.get(ext, "image/png")

        b64 = base64.b64encode(image_bytes).decode("ascii")
        data_uri = f"data:{mime};base64,{b64}"

        client = AzureOpenAI(
            api_key=settings.azure_openai_api_key,
            azure_endpoint=settings.azure_openai_endpoint,
            api_version=settings.azure_openai_api_version,
        )
        resp = client.chat.completions.create(
            model=settings.azure_openai_deployment,
            messages=[
                {"role": "system", "content": (
                    "You are an OCR assistant. Extract ALL text visible in the "
                    "image. Preserve formatting, tables, and structure as much as "
                    "possible. If the image contains no readable text, reply with "
                    "exactly: [NO TEXT FOUND]"
                )},
                {"role": "user", "content": [
                    {"type": "text", "text": f"Extract all text from this image ({filename}):"},
                    {"type": "image_url", "image_url": {"url": data_uri, "detail": "high"}},
                ]},
            ],
            max_tokens=4096,
            temperature=0,
        )
        text = resp.choices[0].message.content.strip()
        if text == "[NO TEXT FOUND]":
            logger.info("GPT-4o vision: no text found in %s", filename)
            return ""
        logger.info("GPT-4o vision OCR extracted %d chars from %s", len(text), filename)
        return text
    except Exception as exc:
        logger.warning("GPT-4o vision OCR failed for %s: %s — trying Tesseract fallback", filename, exc)
        # Tesseract fallback
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(image_bytes))
            return _ocr_image(img)
        except Exception:
            return ""


def _transcribe_audio(audio_bytes: bytes, filename: str = "audio.mp3") -> str:
    """Transcribe audio using Azure OpenAI Whisper, with fallbacks.

    Fallback chain:
      1. Azure OpenAI Whisper deployment
      2. GPT-4o audio input (base64-encoded audio in a chat message)
      3. Local ``whisper`` package (openai-whisper)
    """
    ext = Path(filename).suffix.lower() or ".mp3"

    # ── Attempt 1: Azure OpenAI Whisper ────────────────────────────
    try:
        from openai import AzureOpenAI

        client = AzureOpenAI(
            api_key=settings.azure_openai_api_key,
            azure_endpoint=settings.azure_openai_endpoint,
            api_version=settings.azure_openai_api_version,
        )
        # The API expects a file-like object with a name attribute
        audio_file = io.BytesIO(audio_bytes)
        audio_file.name = filename

        transcript = client.audio.transcriptions.create(
            model=settings.azure_whisper_deployment,
            file=audio_file,
            response_format="text",
        )
        text = transcript.strip() if isinstance(transcript, str) else transcript.text.strip()
        logger.info("Azure Whisper transcribed %d chars from %s", len(text), filename)
        return text
    except Exception as exc:
        logger.warning("Azure Whisper failed for %s: %s — trying GPT-4o audio", filename, exc)

    # ── Attempt 2: GPT-4o with audio input ─────────────────────────
    try:
        from openai import AzureOpenAI

        client = AzureOpenAI(
            api_key=settings.azure_openai_api_key,
            azure_endpoint=settings.azure_openai_endpoint,
            api_version=settings.azure_openai_api_version,
        )
        b64_audio = base64.b64encode(audio_bytes).decode("ascii")
        mime_map = {".mp3": "audio/mp3", ".wav": "audio/wav", ".m4a": "audio/mp4",
                    ".ogg": "audio/ogg", ".flac": "audio/flac", ".webm": "audio/webm"}
        audio_mime = mime_map.get(ext, "audio/mp3")

        resp = client.chat.completions.create(
            model=settings.azure_openai_deployment,
            messages=[
                {"role": "system", "content": (
                    "You are a transcription assistant. Transcribe all spoken "
                    "words in the audio exactly as spoken. If the audio contains "
                    "no speech (only music, tones, or silence), reply with "
                    "exactly: [NO SPEECH DETECTED]"
                )},
                {"role": "user", "content": [
                    {"type": "text", "text": f"Transcribe this audio file ({filename}):"},
                    {"type": "input_audio", "input_audio": {
                        "data": b64_audio,
                        "format": ext.lstrip(".") if ext.lstrip(".") in ("mp3", "wav") else "mp3",
                    }},
                ]},
            ],
            max_tokens=4096,
            temperature=0,
        )
        text = resp.choices[0].message.content.strip()
        if text == "[NO SPEECH DETECTED]":
            logger.info("GPT-4o audio: no speech in %s", filename)
            return ""
        logger.info("GPT-4o audio transcribed %d chars from %s", len(text), filename)
        return text
    except Exception as exc:
        logger.warning("GPT-4o audio failed for %s: %s — trying local whisper", filename, exc)

    # ── Attempt 3: Local whisper package ───────────────────────────
    try:
        import whisper

        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(audio_bytes)
            tmp_path = tmp.name
        try:
            model = whisper.load_model("base")
            result = model.transcribe(tmp_path)
            text = result.get("text", "").strip()
            logger.info("Local Whisper transcribed %d chars from %s", len(text), filename)
            return text
        finally:
            os.unlink(tmp_path)
    except ImportError:
        logger.warning("Local whisper package not installed. Run: pip install openai-whisper")
    except Exception as exc:
        logger.warning("Local whisper failed for %s: %s", filename, exc)

    return ""


def _process_attachment(content_type: str, payload: bytes, filename: str) -> str:
    """Process a single email attachment — returns extracted text or ''.

    Dispatches to GPT-4o vision for images and Whisper for audio files.
    """
    if not settings.process_attachments:
        return ""

    ct = content_type.lower()
    ext = Path(filename).suffix.lower() if filename else ""

    # Image attachments → OCR
    if ct in _IMAGE_MIMES or ext in _IMAGE_EXTS:
        logger.info("Processing image attachment: %s (%s)", filename, ct)
        text = _ocr_image_gpt4o(payload, filename)
        if text:
            return f"\n[ATTACHMENT: {filename} — Image OCR]\n{text}\n[/ATTACHMENT]"
        return ""

    # Audio attachments → Whisper
    if ct in _AUDIO_MIMES or ext in _AUDIO_EXTS:
        logger.info("Processing audio attachment: %s (%s)", filename, ct)
        text = _transcribe_audio(payload, filename)
        if text:
            return f"\n[ATTACHMENT: {filename} — Audio Transcript]\n{text}\n[/ATTACHMENT]"
        return ""

    # PDF attachments → extract text
    if ct == "application/pdf" or ext == ".pdf":
        logger.info("Processing PDF attachment: %s", filename)
        try:
            reader = PdfReader(io.BytesIO(payload))
            pages_text = []
            for page in reader.pages:
                t = page.extract_text() or ""
                if t.strip():
                    pages_text.append(t)
            if pages_text:
                combined = "\n".join(pages_text)
                return f"\n[ATTACHMENT: {filename} — PDF Text]\n{combined}\n[/ATTACHMENT]"
        except Exception as exc:
            logger.warning("PDF attachment extraction failed for %s: %s", filename, exc)
        return ""

    logger.debug("Skipping unsupported attachment type: %s (%s)", filename, ct)
    return ""


# ─────────────────────────────────────────────────────────────────────
#  1. Plain text loader
# ─────────────────────────────────────────────────────────────────────
def load_texts(doc_source: str) -> list[Document]:
    """Load .txt files from a local directory or blob prefix."""
    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".txt"}):
        content = file_bytes.decode("utf-8", errors="replace")
        if not content.strip():
            continue
        docs.append(Document(
            page_content=content,
            metadata={"source": filename, "page": 1, "file_type": "txt"},
        ))
    logger.info("Loaded %d .txt files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  2. PDF loader (with OCR fallback for scanned pages)
# ─────────────────────────────────────────────────────────────────────
def load_pdfs(doc_source: str) -> list[Document]:
    """Load .pdf files; falls back to OCR for pages with no extractable text."""
    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".pdf"}):
        reader = PdfReader(io.BytesIO(file_bytes))
        for page_num, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            # OCR fallback for scanned pages
            if not text.strip():
                try:
                    from PIL import Image
                    for img_obj in page.images:
                        pil_img = Image.open(io.BytesIO(img_obj.data))
                        ocr_text = _ocr_image(pil_img)
                        if ocr_text:
                            text += "\n" + ocr_text
                except Exception as exc:
                    logger.debug("PDF OCR fallback skipped for %s p.%d: %s", filename, page_num, exc)
            if not text.strip():
                continue
            docs.append(Document(
                page_content=text,
                metadata={"source": filename, "page": page_num,
                           "total_pages": len(reader.pages), "file_type": "pdf"},
            ))
    logger.info("Loaded %d PDF pages from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  3. Word (.docx) loader
# ─────────────────────────────────────────────────────────────────────
def load_docx(doc_source: str) -> list[Document]:
    """Load .docx files — extracts paragraphs and tables."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        logger.warning("python-docx not installed — skipping .docx files. Run: pip install python-docx")
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".docx"}):
        try:
            word_doc = DocxDocument(io.BytesIO(file_bytes))
            paragraphs: list[str] = []

            # Extract paragraphs
            for para in word_doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # Extract tables
            for table in word_doc.tables:
                table_rows: list[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))
                if table_rows:
                    paragraphs.append("\n[TABLE]\n" + "\n".join(table_rows) + "\n[/TABLE]")

            content = "\n".join(paragraphs)
            if not content.strip():
                continue

            docs.append(Document(
                page_content=content,
                metadata={"source": filename, "page": 1, "file_type": "docx"},
            ))
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .docx files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  4. Image loader (OCR via Tesseract)
# ─────────────────────────────────────────────────────────────────────
def load_images(doc_source: str) -> list[Document]:
    """Load image files and extract text via Tesseract OCR."""
    try:
        from PIL import Image
    except ImportError:
        logger.warning("Pillow not installed — skipping images. Run: pip install Pillow")
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions=_IMAGE_EXTS):
        try:
            img = Image.open(io.BytesIO(file_bytes))
            text = _ocr_image(img)
            if not text:
                logger.warning("No text extracted from image %s", filename)
                continue
            docs.append(Document(
                page_content=text,
                metadata={"source": filename, "page": 1, "file_type": "image",
                           "image_size": f"{img.width}x{img.height}"},
            ))
        except Exception as exc:
            logger.error("Failed to process image %s: %s", filename, exc)

    logger.info("Loaded %d images (OCR) from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  5. Email (.eml) loader
# ─────────────────────────────────────────────────────────────────────
def _extract_email_text(msg: email.message.EmailMessage) -> tuple[str, list[str], list[tuple[str, str]]]:
    """Extract the plain-text body and process attachments from an email.

    Returns:
        (full_text, attachment_names, attachment_docs) where full_text includes any text
        extracted from image/audio/PDF attachments, and attachment_docs is a list of
        (filename, extracted_text) tuples for creating separate attachment Documents.
    """
    body_parts: list[str] = []
    attachment_names: list[str] = []
    attachment_docs: list[tuple[str, str]] = []  # (filename, text)

    # Headers
    headers = []
    for h in ("From", "To", "Cc", "Date", "Subject"):
        val = msg.get(h)
        if val:
            headers.append(f"{h}: {val}")
    if headers:
        body_parts.append("\n".join(headers))

    # Body + attachments
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            fn = part.get_filename()

            if fn:
                # This is an attachment
                attachment_names.append(fn)
                payload = part.get_payload(decode=True)
                if payload:
                    att_text = _process_attachment(ct, payload, fn)
                    if att_text:
                        body_parts.append(att_text)
                        # Also store for separate attachment Document
                        attachment_docs.append((fn, att_text))
                continue

            if ct == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    body_parts.append(payload.decode("utf-8", errors="replace"))
            elif ct == "text/html":
                # Basic HTML stripping as fallback
                payload = part.get_payload(decode=True)
                if payload:
                    import re
                    html = payload.decode("utf-8", errors="replace")
                    text = re.sub(r"<[^>]+>", " ", html)
                    text = re.sub(r"\s+", " ", text).strip()
                    if text:
                        body_parts.append(text)
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            body_parts.append(payload.decode("utf-8", errors="replace"))

    return "\n\n".join(body_parts), attachment_names, attachment_docs


def load_eml(doc_source: str) -> list[Document]:
    """Load .eml email files with full attachment processing.

    Creates separate attachment Documents with ``parent_message_id`` set
    so the retriever can link attachments back to parent emails.
    """
    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".eml"}):
        try:
            msg = email.message_from_bytes(file_bytes, policy=policy.default)
            text, attachments, attachment_docs = _extract_email_text(msg)
            if not text.strip():
                continue

            message_id = msg.get("Message-ID", "") or ""
            in_reply_to = msg.get("In-Reply-To", "") or ""
            references = msg.get("References", "") or ""
            thread_id = in_reply_to or (references.split()[0] if references.strip() else "")
            clean_msg_id = message_id.strip("<>")

            docs.append(Document(
                page_content=text,
                metadata={
                    "source": filename, "page": 1, "file_type": "eml",
                    "email_from": msg.get("From", ""),
                    "email_to": msg.get("To", ""),
                    "email_subject": msg.get("Subject", ""),
                    "email_date": msg.get("Date", ""),
                    "message_id": clean_msg_id,
                    "thread_id": thread_id.strip("<>"),
                    "attachments": ", ".join(attachments) if attachments else "",
                    "attachment_names": ", ".join(attachments) if attachments else "",
                },
            ))

            # Create separate attachment Documents with parent_message_id linkage
            if settings.attachment_context_link and attachment_docs:
                for att_name, att_text in attachment_docs:
                    docs.append(Document(
                        page_content=att_text,
                        metadata={
                            "source": f"{filename}::{att_name}",
                            "page": 1,
                            "file_type": "attachment",
                            "email_from": msg.get("From", ""),
                            "email_to": msg.get("To", ""),
                            "email_subject": msg.get("Subject", ""),
                            "email_date": msg.get("Date", ""),
                            "parent_message_id": clean_msg_id,
                            "thread_id": thread_id.strip("<>"),
                            "attachment_names": att_name,
                        },
                    ))
                logger.info("Created %d separate attachment docs for %s",
                            len(attachment_docs), filename)
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .eml files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  6. Outlook message (.msg) loader
# ─────────────────────────────────────────────────────────────────────
def load_msg(doc_source: str) -> list[Document]:
    """Load .msg Outlook message files with attachment processing."""
    try:
        import extract_msg
    except ImportError:
        logger.warning("extract-msg not installed — skipping .msg files. Run: pip install extract-msg")
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".msg"}):
        tmp_path = None
        try:
            # extract_msg requires a filesystem path — write bytes to temp file
            with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            msg = extract_msg.Message(tmp_path)
            parts: list[str] = []

            # Headers
            headers = []
            if msg.sender:    headers.append(f"From: {msg.sender}")
            if msg.to:        headers.append(f"To: {msg.to}")
            if msg.cc:        headers.append(f"Cc: {msg.cc}")
            if msg.date:      headers.append(f"Date: {msg.date}")
            if msg.subject:   headers.append(f"Subject: {msg.subject}")
            if headers:
                parts.append("\n".join(headers))

            # Body
            body = msg.body or ""
            if body.strip():
                parts.append(body)

            # Process attachments (image OCR, audio transcription, PDF)
            att_names: list[str] = []
            for att in (msg.attachments or []):
                fn = att.longFilename or att.shortFilename or "unnamed"
                att_names.append(fn)
                if settings.process_attachments:
                    try:
                        att_data = att.data
                        if att_data:
                            # Guess content type from extension
                            ext = Path(fn).suffix.lower()
                            mime_map = {
                                ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                                ".gif": "image/gif", ".bmp": "image/bmp", ".tiff": "image/tiff",
                                ".mp3": "audio/mpeg", ".wav": "audio/wav", ".m4a": "audio/m4a",
                                ".ogg": "audio/ogg", ".flac": "audio/flac", ".webm": "audio/webm",
                                ".pdf": "application/pdf",
                            }
                            ct = mime_map.get(ext, "application/octet-stream")
                            att_text = _process_attachment(ct, att_data, fn)
                            if att_text:
                                parts.append(att_text)
                    except Exception as exc:
                        logger.warning("Failed to process .msg attachment %s: %s", fn, exc)

            content = "\n\n".join(parts)
            if not content.strip():
                msg.close()
                continue

            docs.append(Document(
                page_content=content,
                metadata={
                    "source": filename, "page": 1, "file_type": "msg",
                    "email_from": msg.sender or "",
                    "email_to": msg.to or "",
                    "email_subject": msg.subject or "",
                    "email_date": str(msg.date or ""),
                    "attachments": ", ".join(att_names) if att_names else "",
                },
            ))
            msg.close()
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)

    logger.info("Loaded %d .msg files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  7. Outlook PST loader (Windows only — uses Outlook COM)
# ─────────────────────────────────────────────────────────────────────

# Module-level accumulator: populated by _extract_pst_folder, consumed by ingest().
_pst_metadata_rows: list[dict] = []


def _derive_message_id(item, pst_name: str, index: int) -> str:
    """Derive a stable MessageId with fallback chain:
    1) Internet Message-ID header (PR_INTERNET_MESSAGE_ID)
    2) Outlook EntryID
    3) Deterministic hash from subject + sender + date
    """
    # Try internet message id via MAPI property
    try:
        pr_internet_msg_id = "http://schemas.microsoft.com/mapi/proptag/0x1035001F"
        internet_id = item.PropertyAccessor.GetProperty(pr_internet_msg_id)
        if internet_id and internet_id.strip():
            return internet_id.strip()
    except Exception:
        pass

    # Fallback: EntryID
    entry_id = getattr(item, "EntryID", "")
    if entry_id:
        return entry_id

    # Fallback: deterministic hash
    import hashlib
    raw = f"{pst_name}|{getattr(item, 'Subject', '')}|{getattr(item, 'SenderName', '')}|{getattr(item, 'ReceivedTime', '')}|{index}"
    return hashlib.sha256(raw.encode()).hexdigest()


def _derive_sender_email(item) -> str:
    """Extract sender email address with fallback to empty string."""
    try:
        return getattr(item, "SenderEmailAddress", "") or ""
    except Exception:
        return ""


def load_pst(doc_source: str) -> list[Document]:
    """Load .pst Outlook data files via Outlook COM automation (Windows only).

    Requires Microsoft Outlook installed on the machine.
    When reading from blob storage the PST file is materialized as a local
    temp file because Outlook COM requires a real filesystem path.
    """
    if platform.system() != "Windows":
        logger.warning("PST loading requires Windows + Outlook. Skipping .pst files.")
        return []

    try:
        import win32com.client
    except ImportError:
        logger.warning("pywin32 not installed — skipping .pst files. Run: pip install pywin32")
        return []

    docs: list[Document] = []
    try:
        outlook = win32com.client.Dispatch("Outlook.Application").GetNamespace("MAPI")
    except Exception as exc:
        logger.error("Cannot start Outlook COM — is Outlook installed? %s", exc)
        return []

    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".pst"}):
        tmp_path = None
        try:
            # Outlook COM requires a real filesystem path — write to temp file
            with tempfile.NamedTemporaryFile(suffix=".pst", delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            pst_abs = str(Path(tmp_path).resolve())

            # Add the PST as a store
            outlook.AddStore(pst_abs)
            # Find the store we just added (it appears as the last folder)
            store = None
            for folder in outlook.Folders:
                # Match by checking all folders; the added PST is typically last
                store = folder
            if store is None:
                logger.warning("Could not find added PST store for %s", filename)
                continue

            # Recursively extract emails from all sub-folders
            _extract_pst_folder(store, filename, docs)

            # Remove the PST store
            outlook.RemoveStore(store)
        except Exception as exc:
            logger.error("Failed to process PST %s: %s", filename, exc)
        finally:
            if tmp_path and os.path.exists(tmp_path):
                # Outlook COM may hold the file for a moment - retry with delay
                import time
                for attempt in range(5):
                    try:
                        os.unlink(tmp_path)
                        break
                    except PermissionError:
                        if attempt < 4:
                            time.sleep(0.5)
                        else:
                            logger.warning("Could not delete temp PST file %s (in use)", tmp_path)

    logger.info("Loaded %d emails from .pst files in %s", len(docs), doc_source)
    return docs


def _extract_pst_folder(folder, pst_name: str, docs: list[Document], max_emails: int = 5000) -> None:
    """Recursively extract emails from an Outlook folder tree."""
    try:
        items = folder.Items
        for i in range(1, min(items.Count + 1, max_emails + 1)):
            try:
                item = items.Item(i)
                if item.Class == 43:  # olMail
                    parts: list[str] = []
                    parts.append(f"From: {getattr(item, 'SenderName', '')}")
                    parts.append(f"To: {getattr(item, 'To', '')}")
                    if getattr(item, 'CC', ''):
                        parts.append(f"Cc: {item.CC}")
                    parts.append(f"Date: {getattr(item, 'ReceivedTime', '')}")
                    parts.append(f"Subject: {getattr(item, 'Subject', '')}")
                    body = getattr(item, "Body", "") or ""
                    # Strip leading email headers from Outlook Body to get pure body
                    _body_lines = body.splitlines()
                    _body_start = 0
                    for _bl_idx, _bl in enumerate(_body_lines):
                        stripped = _bl.strip().rstrip('\r')
                        if stripped and not re.match(
                            r'^(From|To|Cc|CC|Bcc|BCC|Date|Subject|Sent|Importance|Attachments)\s*:', stripped
                        ):
                            _body_start = _bl_idx
                            break
                    pure_body = "\n".join(_body_lines[_body_start:]).strip()
                    if body.strip():
                        parts.append(f"\n{body}")

                    # Process attachments (image OCR, audio transcription)
                    att_names = []
                    for j in range(1, item.Attachments.Count + 1):
                        att = item.Attachments.Item(j)
                        fn = att.FileName
                        att_names.append(fn)

                        if settings.process_attachments:
                            ext = Path(fn).suffix.lower()
                            if ext in _IMAGE_EXTS or ext in _AUDIO_EXTS or ext == ".pdf":
                                try:
                                    # Save attachment to temp file, read bytes
                                    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                                        att.SaveAsFile(tmp.name)
                                        tmp_path = tmp.name
                                    att_data = Path(tmp_path).read_bytes()
                                    os.unlink(tmp_path)

                                    mime_map = {
                                        ".png": "image/png", ".jpg": "image/jpeg",
                                        ".jpeg": "image/jpeg", ".gif": "image/gif",
                                        ".bmp": "image/bmp", ".tiff": "image/tiff",
                                        ".mp3": "audio/mpeg", ".wav": "audio/wav",
                                        ".m4a": "audio/m4a", ".ogg": "audio/ogg",
                                        ".pdf": "application/pdf",
                                    }
                                    ct = mime_map.get(ext, "application/octet-stream")
                                    att_text = _process_attachment(ct, att_data, fn)
                                    if att_text:
                                        parts.append(att_text)
                                except Exception as exc:
                                    logger.warning("PST attachment processing failed for %s: %s", fn, exc)

                    content = "\n".join(parts)
                    if content.strip():
                        # Extract message identifiers for threading
                        msg_id = ""
                        conversation_id = ""
                        try:
                            msg_id = getattr(item, "EntryID", "") or ""
                            conversation_id = getattr(item, "ConversationID", "") or getattr(item, "ConversationTopic", "") or ""
                        except Exception:
                            pass

                        docs.append(Document(
                            page_content=content,
                            metadata={
                                "source": pst_name,
                                "page": i,
                                "file_type": "pst",
                                "email_from": getattr(item, "SenderName", ""),
                                "email_to": getattr(item, "To", ""),
                                "email_subject": getattr(item, "Subject", ""),
                                "email_date": str(getattr(item, "ReceivedTime", "")),
                                "message_id": msg_id,
                                "thread_id": conversation_id,
                                "folder": folder.Name,
                                "attachments": ", ".join(att_names) if att_names else "",
                                "attachment_names": ", ".join(att_names) if att_names else "",
                                "body_only": pure_body,
                            },
                        ))

                        # ── Collect PST email metadata row ──────────
                        try:
                            stable_msg_id = _derive_message_id(item, pst_name, i)
                            html_body = ""
                            try:
                                html_body = getattr(item, "HTMLBody", "") or ""
                            except Exception:
                                pass
                            _pst_metadata_rows.append({
                                "PstFileId": pst_name,
                                "Subject": getattr(item, "Subject", "") or "",
                                "SenderName": getattr(item, "SenderName", "") or "",
                                "SenderEmail": _derive_sender_email(item),
                                "RecipientTo": getattr(item, "To", "") or "",
                                "RecipientCc": getattr(item, "CC", "") or "",
                                "SentDate": str(getattr(item, "ReceivedTime", "")),
                                "BodyText": body,
                                "BodyHtml": html_body,
                                "HasAttachments": len(att_names) > 0,
                                "FolderPath": folder.Name,
                                "MessageId": stable_msg_id,
                            })
                        except Exception as meta_exc:
                            logger.warning("Failed to collect PST metadata for email %d in %s: %s", i, pst_name, meta_exc)

                        # Create separate attachment Documents with parent_message_id
                        if settings.attachment_context_link and att_names:
                            for att_fn in att_names:
                                # Find matching attachment text in parts
                                att_marker = f"[ATTACHMENT: {att_fn}"
                                for part in parts:
                                    if att_marker in part:
                                        docs.append(Document(
                                            page_content=part,
                                            metadata={
                                                "source": f"{pst_name}::{att_fn}",
                                                "page": i,
                                                "file_type": "attachment",
                                                "email_from": getattr(item, "SenderName", ""),
                                                "email_to": getattr(item, "To", ""),
                                                "email_subject": getattr(item, "Subject", ""),
                                                "email_date": str(getattr(item, "ReceivedTime", "")),
                                                "parent_message_id": msg_id,
                                                "thread_id": conversation_id,
                                                "attachment_names": att_fn,
                                            },
                                        ))
                                        break
            except Exception:
                continue
    except Exception:
        pass

    # Recurse into sub-folders
    try:
        for j in range(1, folder.Folders.Count + 1):
            _extract_pst_folder(folder.Folders.Item(j), pst_name, docs, max_emails)
    except Exception:
        pass


# ─────────────────────────────────────────────────────────────────────
#  8. Purview eDiscovery export loader (JSON / CSV)
# ─────────────────────────────────────────────────────────────────────
def _flatten_json_for_text(value: object, depth: int = 0) -> str:
    """Convert JSON-like data into readable text for chunking/embedding."""
    if depth >= _JSON_MAX_DEPTH:
        return "[truncated]"
    if isinstance(value, dict):
        lines: list[str] = []
        for key, val in value.items():
            if isinstance(val, (dict, list)):
                child = _flatten_json_for_text(val, depth + 1)
                if child:
                    lines.append(f"{key}:")
                    lines.append(child)
            else:
                lines.append(f"{key}: {val}")
        return "\n".join(lines)
    if isinstance(value, list):
        return "\n".join(_flatten_json_for_text(item, depth + 1) for item in value)
    return str(value)


def load_generic_json(doc_source: str) -> list[Document]:
    """Load non-Purview JSON files as readable text."""
    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".json"}):
        # Items*.json is handled by the Purview-specific loader below.
        if filename.lower().startswith("items"):
            continue
        try:
            data = json.loads(file_bytes.decode("utf-8", errors="replace"))
        except Exception as exc:
            logger.error("Failed to parse JSON %s: %s", filename, exc)
            continue

        text = _flatten_json_for_text(data).strip()
        if not text:
            continue

        docs.append(
            Document(
                page_content=text,
                metadata={"source": filename, "page": 1, "file_type": "json"},
            )
        )

    logger.info("Loaded %d generic JSON files from %s", len(docs), doc_source)
    return docs


def _html_to_text(html_text: str) -> str:
    """Best-effort HTML to plain text without external dependencies."""
    cleaned = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", html_text)
    cleaned = re.sub(r"(?i)<br\s*/?>", "\n", cleaned)
    cleaned = re.sub(r"(?i)</p\s*>", "\n", cleaned)
    cleaned = re.sub(r"(?is)<[^>]+>", " ", cleaned)
    cleaned = unescape(cleaned)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def load_html(doc_source: str) -> list[Document]:
    """Load .html/.htm files and convert to plain text."""
    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".html", ".htm"}):
        html_text = file_bytes.decode("utf-8", errors="replace")
        text = _html_to_text(html_text)
        if not text:
            continue
        docs.append(
            Document(
                page_content=text,
                metadata={"source": filename, "page": 1, "file_type": "html"},
            )
        )
    logger.info("Loaded %d HTML files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  9. Excel spreadsheet (.xlsx) loader
# ─────────────────────────────────────────────────────────────────────


def load_xlsx(doc_source: str) -> list[Document]:
    """Load .xlsx files — extracts each sheet as a table."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed — skipping .xlsx files. Run: pip install openpyxl")
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".xlsx"}):
        try:
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
                ws = wb[sheet_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if not rows:
                    continue
                content = f"Sheet: {sheet_name}\n\n[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]"
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": filename, "page": sheet_idx,
                        "file_type": "xlsx", "sheet_name": sheet_name,
                        "row_count": len(rows),
                    },
                ))
            wb.close()
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .xlsx sheets from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  10. Legacy Excel spreadsheet (.xls) loader
# ─────────────────────────────────────────────────────────────────────


def load_xls(doc_source: str) -> list[Document]:
    """Load .xls files (legacy Excel 97-2003) — extracts each sheet as a table."""
    try:
        import xlrd
    except ImportError:
        logger.warning("xlrd not installed — skipping .xls files. Run: pip install xlrd>=2.0.1")
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".xls"}):
        try:
            wb = xlrd.open_workbook(file_contents=file_bytes)
            for sheet_idx in range(wb.nsheets):
                ws = wb.sheet_by_index(sheet_idx)
                rows: list[str] = []
                for row_num in range(ws.nrows):
                    cells = [str(ws.cell_value(row_num, col)) for col in range(ws.ncols)]
                    if any(c.strip() for c in cells):
                        rows.append(" | ".join(cells))
                if not rows:
                    continue
                content = f"Sheet: {ws.name}\n\n[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]"
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": filename, "page": sheet_idx + 1,
                        "file_type": "xls", "sheet_name": ws.name,
                        "row_count": len(rows),
                    },
                ))
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .xls sheets from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  11. PowerPoint (.pptx) loader
# ─────────────────────────────────────────────────────────────────────


def load_pptx(doc_source: str) -> list[Document]:
    """Load .pptx files — extracts slide text, notes, and table content."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed — skipping .pptx files. Run: pip install python-pptx")
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".pptx"}):
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            slide_count = len(prs.slides)
            for slide_num, slide in enumerate(prs.slides, start=1):
                parts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            parts.append(text)
                    if shape.has_table:
                        table_rows: list[str] = []
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            table_rows.append(" | ".join(cells))
                        if table_rows:
                            parts.append("\n[TABLE]\n" + "\n".join(table_rows) + "\n[/TABLE]")
                # Slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[NOTES]\n{notes}\n[/NOTES]")
                content = "\n\n".join(parts)
                if not content.strip():
                    continue
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": filename, "page": slide_num,
                        "file_type": "pptx", "slide_count": slide_count,
                    },
                ))
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .pptx slides from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  12. Markdown (.md) loader
# ─────────────────────────────────────────────────────────────────────


def load_markdown(doc_source: str) -> list[Document]:
    """Load .md Markdown files as plain text."""
    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".md"}):
        content = file_bytes.decode("utf-8", errors="replace")
        if not content.strip():
            continue
        docs.append(Document(
            page_content=content,
            metadata={"source": filename, "page": 1, "file_type": "markdown"},
        ))
    logger.info("Loaded %d .md files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  13. Rich Text Format (.rtf) loader
# ─────────────────────────────────────────────────────────────────────


def load_rtf(doc_source: str) -> list[Document]:
    """Load .rtf files — strips RTF formatting to extract plain text."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        logger.warning("striprtf not installed — skipping .rtf files. Run: pip install striprtf")
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".rtf"}):
        try:
            raw = file_bytes.decode("utf-8", errors="replace")
            content = rtf_to_text(raw)
            if not content.strip():
                continue
            docs.append(Document(
                page_content=content,
                metadata={"source": filename, "page": 1, "file_type": "rtf"},
            ))
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .rtf files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  14. Legacy Word binary (.doc) loader — best-effort
# ─────────────────────────────────────────────────────────────────────


def load_doc(doc_source: str) -> list[Document]:
    """Load .doc files via olefile — best-effort text extraction.

    Legacy .doc format has no reliable pure-python parser.  This extracts
    the raw text stream from the OLE container.  For full fidelity,
    convert .doc to .docx first.
    """
    try:
        import olefile
    except ImportError:
        logger.warning(
            "olefile not installed — skipping .doc files. "
            "Run: pip install olefile  (or convert .doc to .docx for best results)"
        )
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".doc"}):
        try:
            ole = olefile.OleFileIO(io.BytesIO(file_bytes))
            if ole.exists("WordDocument"):
                stream = ole.openstream("WordDocument").read()
                text = stream.decode("utf-8", errors="ignore")
                text = re.sub(r"[^\x20-\x7E\n\r\t]", "", text)
                text = re.sub(r"\n{3,}", "\n\n", text).strip()
                if text and len(text) > 20:
                    docs.append(Document(
                        page_content=text,
                        metadata={"source": filename, "page": 1, "file_type": "doc"},
                    ))
                else:
                    logger.warning("Minimal text extracted from %s — consider converting to .docx", filename)
            else:
                logger.warning("No WordDocument stream in %s — cannot extract text", filename)
            ole.close()
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .doc files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  15. XML document loader
# ─────────────────────────────────────────────────────────────────────


def load_xml(doc_source: str) -> list[Document]:
    """Load .xml files — extracts all text content from elements."""
    try:
        import defusedxml.ElementTree as ET
    except ImportError:
        logger.warning("defusedxml not installed; falling back to xml.etree (unsafe for untrusted XML)")
        import xml.etree.ElementTree as ET

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".xml"}):
        try:
            tree = ET.parse(io.BytesIO(file_bytes))
            root = tree.getroot()
            root_tag = root.tag.split("}")[-1] if "}" in root.tag else root.tag
            texts: list[str] = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                    texts.append(f"{tag}: {elem.text.strip()}")
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())
            content = "\n".join(texts)
            if not content.strip():
                continue
            docs.append(Document(
                page_content=content,
                metadata={
                    "source": filename, "page": 1,
                    "file_type": "xml", "root_tag": root_tag,
                },
            ))
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .xml files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  16. Generic CSV loader (non-Purview)
# ─────────────────────────────────────────────────────────────────────


def load_generic_csv(doc_source: str) -> list[Document]:
    """Load generic .csv files as tabular data.

    Purview Items_*.csv files are handled by load_purview_csv and skipped here.
    Rows are batched into groups of _CSV_BATCH_SIZE to create manageable chunks.
    """
    import csv as _csv

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".csv"}):
        if filename.lower().startswith("items"):
            continue
        try:
            text_data = file_bytes.decode("utf-8-sig", errors="replace")
            reader = _csv.reader(io.StringIO(text_data))
            all_rows = list(reader)
            if not all_rows:
                continue
            header = all_rows[0]
            header_line = " | ".join(header)
            data_rows = all_rows[1:]
            total_rows = len(data_rows)

            for batch_idx in range(0, max(len(data_rows), 1), _CSV_BATCH_SIZE):
                batch = data_rows[batch_idx: batch_idx + _CSV_BATCH_SIZE]
                if not batch:
                    continue
                row_lines = [" | ".join(r) for r in batch]
                content = f"[TABLE]\n{header_line}\n" + "\n".join(row_lines) + "\n[/TABLE]"
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": filename,
                        "page": (batch_idx // _CSV_BATCH_SIZE) + 1,
                        "file_type": "csv",
                        "row_count": total_rows,
                    },
                ))
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d generic CSV chunks from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  17. Tab-separated values (.tsv) loader
# ─────────────────────────────────────────────────────────────────────


def load_tsv(doc_source: str) -> list[Document]:
    """Load .tsv files as tabular data with the same batching as CSV."""
    import csv as _csv

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".tsv"}):
        try:
            text_data = file_bytes.decode("utf-8-sig", errors="replace")
            reader = _csv.reader(io.StringIO(text_data), delimiter="\t")
            all_rows = list(reader)
            if not all_rows:
                continue
            header = all_rows[0]
            header_line = " | ".join(header)
            data_rows = all_rows[1:]
            total_rows = len(data_rows)

            for batch_idx in range(0, max(len(data_rows), 1), _CSV_BATCH_SIZE):
                batch = data_rows[batch_idx: batch_idx + _CSV_BATCH_SIZE]
                if not batch:
                    continue
                row_lines = [" | ".join(r) for r in batch]
                content = f"[TABLE]\n{header_line}\n" + "\n".join(row_lines) + "\n[/TABLE]"
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": filename,
                        "page": (batch_idx // _CSV_BATCH_SIZE) + 1,
                        "file_type": "tsv",
                        "row_count": total_rows,
                    },
                ))
        except Exception as exc:
            logger.error("Failed to load %s: %s", filename, exc)

    logger.info("Loaded %d .tsv chunks from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  18. Log file (.log) loader
# ─────────────────────────────────────────────────────────────────────


def load_logs(doc_source: str) -> list[Document]:
    """Load .log files as plain text."""
    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".log"}):
        content = file_bytes.decode("utf-8", errors="replace")
        if not content.strip():
            continue
        docs.append(Document(
            page_content=content,
            metadata={"source": filename, "page": 1, "file_type": "log"},
        ))
    logger.info("Loaded %d .log files from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  19. ZIP archive loader (recursive extraction with safety guards)
# ─────────────────────────────────────────────────────────────────────


def load_zip(doc_source: str, _depth: int = 0) -> list[Document]:
    """Load .zip files by extracting contents and recursively loading.

    Safety guards: max nesting depth, max total bytes, max file count,
    path traversal rejection.
    """
    import shutil
    import zipfile

    if _depth >= _ZIP_MAX_NESTING_DEPTH:
        logger.warning("ZIP nesting depth %d exceeds limit %d — skipping", _depth, _ZIP_MAX_NESTING_DEPTH)
        return []

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".zip"}):
        tmp_dir = None
        try:
            bio = io.BytesIO(file_bytes)
            if not zipfile.is_zipfile(bio):
                logger.warning("Not a valid ZIP file: %s", filename)
                continue

            bio.seek(0)
            zf = zipfile.ZipFile(bio)

            # Safety: check file count
            if len(zf.namelist()) > _ZIP_MAX_FILE_COUNT:
                logger.warning("ZIP %s contains %d files (limit %d) — skipping",
                               filename, len(zf.namelist()), _ZIP_MAX_FILE_COUNT)
                zf.close()
                continue

            # Safety: check total uncompressed size
            total_size = sum(info.file_size for info in zf.infolist())
            if total_size > _ZIP_MAX_TOTAL_BYTES:
                logger.warning("ZIP %s uncompressed size %.1f MB exceeds limit — skipping",
                               filename, total_size / 1024 / 1024)
                zf.close()
                continue

            # Safety: check for path traversal using resolved paths
            bad_path = False
            tmp_dir = tempfile.mkdtemp(prefix="rag_zip_")
            extract_root = Path(tmp_dir).resolve()
            for entry in zf.infolist():
                target = (extract_root / entry.filename).resolve()
                if not str(target).startswith(str(extract_root)):
                    logger.warning("ZIP %s contains path-traversal entry '%s' — skipping archive", filename, entry.filename)
                    bad_path = True
                    break
            if bad_path:
                zf.close()
                import shutil as _shutil
                _shutil.rmtree(tmp_dir, ignore_errors=True)
                tmp_dir = None
                continue

            zf.extractall(tmp_dir)
            zf.close()

            # Recursively load extracted files
            inner_docs = load_documents(tmp_dir, _zip_depth=_depth + 1)

            for doc in inner_docs:
                doc.metadata["archive_source"] = filename
                doc.metadata["source"] = f"{filename}::{doc.metadata.get('source', 'unknown')}"

            docs.extend(inner_docs)
            logger.info("Extracted %d documents from ZIP %s", len(inner_docs), filename)

        except Exception as exc:
            logger.error("Failed to process ZIP %s: %s", filename, exc)
        finally:
            if tmp_dir and os.path.isdir(tmp_dir):
                shutil.rmtree(tmp_dir, ignore_errors=True)

    logger.info("Loaded %d documents from .zip archives in %s", len(docs), doc_source)
    return docs


_PURVIEW_KEY_FIELDS = [
    "Subject/Title", "Sender/Author", "To", "CC", "BCC",
    "Date", "Participants", "Conversation topic", "Themes list",
    "Message kind", "File class", "Custodian", "Attachment names",
    "Email importance", "Compound path", "Dominant theme",
    "Word count", "Extracted text length",
]


def _purview_item_to_text(item: dict) -> str:
    """Build a searchable text block from a Purview metadata record."""
    parts: list[str] = []
    subject = item.get("Subject/Title", "").strip()
    if subject:
        parts.append(f"Subject: {subject}")

    sender = item.get("Sender/Author", "") or item.get("Sender", "")
    if sender:
        parts.append(f"From: {sender}")

    for field in ("To", "CC", "BCC"):
        val = item.get(field, "").strip()
        if val:
            parts.append(f"{field}: {val}")

    date = item.get("Date", "") or item.get("Email date sent", "")
    if date:
        parts.append(f"Date: {date}")

    topic = item.get("Conversation topic", "").strip()
    if topic:
        parts.append(f"Topic: {topic}")

    participants = item.get("Participants", "").strip()
    if participants:
        parts.append(f"Participants: {participants}")

    kind = item.get("Message kind", "").strip()
    if kind:
        parts.append(f"Type: {kind}")

    attachments = item.get("Attachment names", "").strip()
    if attachments:
        parts.append(f"Attachments: {attachments}")

    themes = item.get("Themes list", "").strip()
    if themes:
        parts.append(f"Themes: {themes}")

    importance = item.get("Email importance", "").strip()
    if importance and importance.lower() != "normal":
        parts.append(f"Importance: {importance}")

    return "\n".join(parts)


def load_purview_json(doc_source: str) -> list[Document]:
    """Load Purview eDiscovery JSON export (Items_*.json files)."""
    import json as _json

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".json"}):
        # Only process Purview Items JSON files
        if not filename.lower().startswith("items"):
            logger.debug("Skipping non-Items JSON: %s", filename)
            continue

        try:
            items = _json.loads(file_bytes.decode("utf-8", errors="replace"))
            if not isinstance(items, list):
                logger.warning("Unexpected JSON structure in %s (not an array)", filename)
                continue

            for idx, item in enumerate(items):
                text = _purview_item_to_text(item)
                if not text.strip():
                    continue

                subject = item.get("Subject/Title", "")
                sender = item.get("Sender/Author", "") or item.get("Sender", "")
                msg_kind = item.get("Message kind", "")
                file_class = item.get("File class", "")

                docs.append(Document(
                    page_content=text,
                    metadata={
                        "source": filename,
                        "page": idx + 1,
                        "file_type": "purview-json",
                        "email_from": sender,
                        "email_to": item.get("To", ""),
                        "email_subject": subject,
                        "email_date": item.get("Date", ""),
                        "message_kind": msg_kind,
                        "file_class": file_class,
                        "custodian": item.get("Custodian", ""),
                        "file_id": item.get("File ID", ""),
                    },
                ))

            logger.info("Loaded %d Purview items from %s", len(docs), filename)
        except Exception as exc:
            logger.error("Failed to load Purview JSON %s: %s", filename, exc)

    logger.info("Loaded %d Purview items total from %s", len(docs), doc_source)
    return docs


def load_purview_csv(doc_source: str) -> list[Document]:
    """Load Purview eDiscovery CSV export (Items_*.csv files)."""
    import csv as _csv

    docs: list[Document] = []
    for filename, file_bytes in iter_files_for_source(doc_source, extensions={".csv"}):
        if not filename.lower().startswith("items"):
            logger.debug("Skipping non-Items CSV: %s", filename)
            continue

        try:
            text_data = file_bytes.decode("utf-8-sig", errors="replace")
            reader = _csv.DictReader(io.StringIO(text_data))

            count = 0
            for idx, row in enumerate(reader):
                item_text = _purview_item_to_text(row)
                if not item_text.strip():
                    continue

                subject = row.get("Subject/Title", "")
                sender = row.get("Sender/Author", "") or row.get("Sender", "")

                docs.append(Document(
                    page_content=item_text,
                    metadata={
                        "source": filename,
                        "page": idx + 1,
                        "file_type": "purview-csv",
                        "email_from": sender,
                        "email_to": row.get("To", ""),
                        "email_subject": subject,
                        "email_date": row.get("Date", ""),
                        "message_kind": row.get("Message kind", ""),
                        "file_class": row.get("File class", ""),
                        "custodian": row.get("Custodian", ""),
                        "file_id": row.get("File ID", ""),
                    },
                ))
                count += 1

            logger.info("Loaded %d Purview items from %s", count, filename)
        except Exception as exc:
            logger.error("Failed to load Purview CSV %s: %s", filename, exc)

    logger.info("Loaded %d Purview items total from %s", len(docs), doc_source)
    return docs


# ─────────────────────────────────────────────────────────────────────
#  Master loader — dispatches to all format-specific loaders
# ─────────────────────────────────────────────────────────────────────
def load_documents(doc_source: str, *, _zip_depth: int = 0) -> list[Document]:
    """Load all supported file types from a local directory or blob prefix.

    Supported: .txt, .pdf, .docx, .xlsx, .xls, .pptx, .md, .rtf, .doc,
               .xml, .csv, .tsv, .log, .json, .html/.htm,
               .png/.jpg/.tiff/.bmp (OCR), .eml, .msg, .pst, .zip
    """
    docs: list[Document] = []
    file_counts: dict[str, int] = {}

    # Each loader calls iter_files_for_source internally and handles
    # empty iteration gracefully, so no pre-check needed.
    for label, loader in [
        ("txt",    load_texts),
        ("pdf",    load_pdfs),
        ("docx",   load_docx),
        ("xlsx",   load_xlsx),
        ("xls",    load_xls),
        ("pptx",   load_pptx),
        ("markdown", load_markdown),
        ("rtf",    load_rtf),
        ("doc",    load_doc),
        ("xml",    load_xml),
        ("generic-csv", load_generic_csv),
        ("tsv",    load_tsv),
        ("log",    load_logs),
        ("json",   load_generic_json),
        ("html",   load_html),
        ("images", load_images),
        ("eml",    load_eml),
        ("msg",    load_msg),
        ("pst",    load_pst),
        ("purview-json", load_purview_json),
        ("purview-csv",  load_purview_csv),
    ]:
        loaded = loader(doc_source)
        if loaded:
            file_counts[label] = len(loaded)
            docs.extend(loaded)

    # ZIP loader receives the current depth for thread-safe nesting control
    zip_loaded = load_zip(doc_source, _depth=_zip_depth)
    if zip_loaded:
        file_counts["zip"] = len(zip_loaded)
        docs.extend(zip_loaded)

    summary = " | ".join(f"{ext}: {cnt}" for ext, cnt in file_counts.items() if cnt)
    logger.info("Total documents loaded: %d  (%s)", len(docs), summary or "none")
    return docs


def chunk_documents(docs: list[Document], chunk_size: int = settings.chunk_size, chunk_overlap: int = settings.chunk_overlap) -> list[Document]:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap, separators=["\n\n", "\n", ". ", " ", ""])
    chunks: list[Document] = []
    for doc in docs:
        splits = splitter.split_text(doc.page_content)
        for idx, text in enumerate(splits):
            chunks.append(Document(page_content=text, metadata={**doc.metadata, "chunk_index": idx}))
    logger.info("Split into %d chunks", len(chunks))
    return chunks


def _get_embedding_model():
    """Return the Cohere embedding model (matching retriever)."""
    try:
        from langchain_cohere import CohereEmbeddings
        return CohereEmbeddings(
            model=settings.cohere_embed_model,
            cohere_api_key=settings.cohere_api_key,
        )
    except Exception:
        # Fallback to HuggingFace (WARNING: 384-dim, won't match Cohere 1024-dim index!)
        logger.warning("Cohere embeddings unavailable, falling back to HuggingFace (dim mismatch risk!)")
        from langchain_community.embeddings import HuggingFaceEmbeddings
        return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")


def embed_chunks(chunks: list[Document], batch_size: int = 48, progress_cb=None) -> list[dict]:
    """Embed chunks in batches with rate-limit retry."""
    import time

    model = _get_embedding_model()
    texts = [c.page_content for c in chunks]
    all_embeddings: list[list[float]] = []

    for i in range(0, len(texts), batch_size):
        batch = texts[i : i + batch_size]
        batch_num = i // batch_size + 1
        total_batches = (len(texts) + batch_size - 1) // batch_size

        for attempt in range(5):
            try:
                embs = model.embed_documents(batch)
                all_embeddings.extend(embs)
                logger.info("Embedded batch %d/%d (%d chunks)", batch_num, total_batches, len(batch))
                if progress_cb:
                    progress_cb(batch_num, total_batches)
                break
            except Exception as exc:
                if "429" in str(exc) or "rate" in str(exc).lower():
                    wait = 2 ** attempt * 5  # 5, 10, 20, 40, 80 seconds
                    logger.warning("Rate limited on batch %d, retrying in %ds...", batch_num, wait)
                    time.sleep(wait)
                else:
                    raise
        else:
            raise RuntimeError(f"Failed to embed batch {batch_num} after 5 retries")

        # Small delay between batches to stay under rate limit
        if i + batch_size < len(texts):
            time.sleep(2)

    rows = []
    for chunk, emb in zip(chunks, all_embeddings):
        meta = chunk.metadata
        # Separate known DB columns from extra metadata
        _db_cols = {
            "source", "page", "chunk_index", "doc_type", "entity_name",
            "effective_date", "message_id", "thread_id", "email_subject",
            "email_from", "email_to", "email_date", "chunk_type",
            "attachment_names", "parent_message_id", "total_chunks",
            "body_only",
        }
        rows.append({
            "content": chunk.page_content, "embedding": emb,
            "source": meta.get("source", "unknown"), "page": meta.get("page"),
            "chunk_index": meta.get("chunk_index"), "doc_type": meta.get("doc_type"),
            "entity_name": meta.get("entity_name"), "effective_date": meta.get("effective_date"),
            # Email-specific metadata
            "message_id": meta.get("message_id", ""),
            "thread_id": meta.get("thread_id", ""),
            "email_subject": meta.get("email_subject", ""),
            "email_from": meta.get("email_from", ""),
            "email_to": meta.get("email_to", ""),
            "email_date": meta.get("email_date", ""),
            "chunk_type": meta.get("chunk_type", ""),
            "attachment_names": meta.get("attachment_names") or meta.get("attachments", ""),
            "parent_message_id": meta.get("parent_message_id", ""),
            "total_chunks": meta.get("total_chunks"),
            "body_only": meta.get("body_only"),
            "metadata_extra": {k: v for k, v in meta.items() if k not in _db_cols},
        })
    return rows


def ingest(
    doc_source: str,
    *,
    case_id: str,
    doc_type: str | None = None,
    entity_name: str | None = None,
    progress_cb=None,
) -> tuple[int, IngestionReport]:
    """Ingest documents into the **per-case** isolated database.

    ``doc_source`` is either a local directory path or a blob prefix
    (when Azure Blob Storage is configured).  ``case_id`` is mandatory
    — it determines which database receives the chunks.

    ``progress_cb`` is an optional ``(pct: float, label: str) -> None``
    callback invoked at each pipeline stage so callers (e.g. UI) can
    display a live progress bar.

    Returns ``(chunks_inserted, report)`` where *report* tracks
    per-file success/failure for auditability.

    Guardrail: when ``settings.require_blob_source`` is True, raises
    RuntimeError if Azure Blob Storage is not configured.
    """
    from app.blob_storage import is_blob_mode

    def _progress(pct: float, label: str) -> None:
        if progress_cb:
            progress_cb(pct, label)

    # ── Guardrail: reject non-blob sources when required ─────────
    if settings.require_blob_source and not is_blob_mode():
        raise RuntimeError(
            "Ingestion BLOCKED: require_blob_source=True but Azure Blob "
            "Storage is not configured. Set AZURE_STORAGE_SAS_URL or "
            "AZURE_STORAGE_CONNECTION_STRING, or set REQUIRE_BLOB_SOURCE=false."
        )

    _progress(0.02, "Initializing database...")
    engine = init_db(case_id)

    # ── Load ─────────────────────────────────────────────────────
    _progress(0.05, "Loading documents...")
    report = IngestionReport()
    _pst_metadata_rows.clear()  # reset before loading
    raw_docs = load_documents(doc_source)
    report.total_files_found = len(raw_docs)
    report.files_loaded = len(raw_docs)
    logger.info("[metrics] files_read=%d from source=%s", len(raw_docs), doc_source)

    # ── Persist PST email metadata (feature-gated to demo-pst) ──
    if case_id == "demo-pst" and _pst_metadata_rows:
        meta_rows = list(_pst_metadata_rows)  # snapshot before clearing
        _pst_metadata_rows.clear()
        for row in meta_rows:
            row["case_id"] = case_id
        meta_count = upsert_pst_email_metadata(engine, meta_rows)
        logger.info(
            "[pst-metadata] Upserted %d email metadata rows for case '%s'",
            meta_count, case_id,
        )

    if not raw_docs:
        logger.warning("No documents loaded from %s — nothing to ingest.", doc_source)
        _progress(1.0, "No documents found")
        return 0, report

    _progress(0.10, f"Loaded {len(raw_docs)} documents")

    if doc_type or entity_name:
        for d in raw_docs:
            if doc_type: d.metadata["doc_type"] = doc_type
            if entity_name: d.metadata["entity_name"] = entity_name

    # ── Chunk ────────────────────────────────────────────────────
    _progress(0.12, "Chunking documents...")
    if settings.use_structure_aware_chunking:
        from app.chunking import chunk_documents_structured
        chunks = chunk_documents_structured(raw_docs)
    else:
        chunks = chunk_documents(raw_docs)
    skipped = len(raw_docs) - len(chunks) if len(chunks) < len(raw_docs) else 0
    logger.info(
        "[metrics] docs=%d  chunks_created=%d  chunks_skipped=%d  "
        "chunk_size=%d  chunk_overlap=%d",
        len(raw_docs), len(chunks), skipped,
        settings.chunk_size, settings.chunk_overlap,
    )
    _progress(0.15, f"Created {len(chunks)} chunks")

    # ── PII redaction (before embedding) ──────────────────────────
    if settings.pii_redaction_enabled:
        _progress(0.16, "Running PII redaction...")
        from app.pii_redaction import redact_document_chunks
        chunks, pii_stats = redact_document_chunks(chunks, mode=settings.pii_redaction_mode)
        if pii_stats.get("total_detections", 0) > 0:
            logger.info(
                "[pii] Redacted %d PII instances across %d/%d chunks in case '%s'",
                pii_stats["total_detections"],
                pii_stats.get("chunks_with_pii", 0),
                len(chunks),
                case_id,
            )
            try:
                from app.audit_log import audit_log
                audit_log.log_pii_event(
                    action="redacted",
                    case_id=case_id,
                    details=pii_stats,
                )
            except Exception:
                pass
    _progress(0.20, "Embedding chunks...")

    # ── Embed ─────────────────────────────────────────────────────
    def _embed_progress(batch_num: int, total_batches: int) -> None:
        # Map batch progress into 0.20 → 0.70 range
        pct = 0.20 + (batch_num / total_batches) * 0.50
        _progress(pct, f"Embedding batch {batch_num}/{total_batches}...")

    rows = embed_chunks(chunks, progress_cb=_embed_progress)
    _progress(0.70, "Embedding complete")

    # ── Enrich (optional, LLM-based) ──────────────────────────────
    if settings.enable_metadata_enrichment:
        _progress(0.72, "Enriching metadata...")
        from app.enrichment import enrich_chunks_batch
        rows = enrich_chunks_batch(rows, batch_size=settings.enrichment_batch_size)
    _progress(0.85, "Storing chunks in database...")

    # ── Store ─────────────────────────────────────────────────────
    inserted = upsert_chunks(engine, rows)
    logger.info(
        "[metrics] embeddings_written=%d  case=%s  source=%s",
        inserted, case_id, doc_source,
    )
    _progress(0.93, f"Stored {inserted} chunks")

    # ── Invalidate retrieval cache for this case ──────────
    # New chunks change search results; embedding cache is safe to keep.
    _progress(0.95, "Invalidating caches...")
    from app.cache import cache
    removed = cache.invalidate_case_retrieval_cache(case_id)
    if removed:
        logger.info("[cache] Invalidated %d retrieval cache entries for case '%s'", removed, case_id)

    report.chunks_created = len(chunks)
    report.chunks_stored = inserted
    logger.info("[report] %s", report.summary())
    if report.files_failed > 0:
        logger.warning("[report] Failed files: %s", [f["file"] for f in report.failed_files])

    _progress(1.0, f"Ingestion complete — {inserted} chunks stored")
    return inserted, report


def ingest_case(case_id: str, source_path: str, progress_cb=None) -> dict:
    """Convenience wrapper for job_queue — ingest by case_id + source_path.

    Resolves doc_type / entity_name from the case config, calls ``ingest()``,
    and returns a JSON-serialisable dict suitable as a job result.
    """
    doc_type = None
    entity_name = None
    try:
        from app.cases import get_case
        cfg = get_case(case_id)
        doc_type = cfg.doc_type
        entity_name = cfg.entity_name
    except (KeyError, Exception):
        pass

    count, report = ingest(
        source_path, case_id=case_id,
        doc_type=doc_type, entity_name=entity_name,
        progress_cb=progress_cb,
    )
    return {
        "case_id": case_id,
        "chunks_inserted": count,
        "summary": report.summary(),
        "files_failed": report.files_failed,
        "failed_files": report.failed_files,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Ingest documents into pgvector")
    parser.add_argument("--dir", dest="doc_dir", help="Local directory with documents")
    parser.add_argument("--blob-prefix", dest="blob_prefix", help="Azure Blob Storage prefix (e.g. big-thorium)")
    parser.add_argument("--pdf-dir", dest="pdf_dir", help="(Legacy) Directory with PDFs")
    parser.add_argument("--case", dest="case_id", required=True, help="Case ID for database isolation (e.g. big-thorium)")
    parser.add_argument("--doc-type", default=None, help="Document type (auto-detected from case config if omitted)")
    parser.add_argument("--entity-name", default=None)
    args = parser.parse_args()

    doc_source = args.blob_prefix or args.doc_dir or args.pdf_dir
    if not doc_source:
        parser.error("Provide --dir, --blob-prefix, or --pdf-dir")

    # Auto-fetch doc_type from case config if not specified
    doc_type = args.doc_type
    entity_name = args.entity_name
    if doc_type is None or entity_name is None:
        try:
            from app.cases import get_case
            case_cfg = get_case(args.case_id)
            if doc_type is None:
                doc_type = case_cfg.doc_type
            if entity_name is None:
                entity_name = case_cfg.entity_name
        except KeyError:
            pass  # Case not registered, use whatever was provided

    logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(name)s | %(message)s")
    count, report = ingest(doc_source, case_id=args.case_id, doc_type=doc_type, entity_name=entity_name)
    print(f"\nIngested {count} chunks into per-case database for '{args.case_id}'.")
    print(f"Report: {report.summary()}")
    if report.files_failed > 0:
        print(f"WARNING: {report.files_failed} file(s) failed to load:")
        for f in report.failed_files:
            print(f"  - {f['file']}: {f['error']}")


if __name__ == "__main__":
    main()
